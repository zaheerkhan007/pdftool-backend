"""
Real PDF processing logic. Each function takes file bytes / paths and returns
processed bytes. Keep these framework-agnostic so they can run inside a Celery
worker or directly in a request.
"""
import io
from typing import List

import pikepdf
from pypdf import PdfReader, PdfWriter


def merge_pdfs(files: List[bytes]) -> bytes:
    """Merge multiple PDFs (in order) into one."""
    writer = PdfWriter()
    for data in files:
        reader = PdfReader(io.BytesIO(data))
        for page in reader.pages:
            writer.add_page(page)
    out = io.BytesIO()
    writer.write(out)
    return out.getvalue()


def split_pdf(data: bytes, ranges: List[tuple]) -> List[bytes]:
    """
    Split a PDF into multiple files by page ranges.
    ranges: list of (start, end) 1-indexed inclusive tuples.
    Returns one PDF (bytes) per range.
    """
    reader = PdfReader(io.BytesIO(data))
    total = len(reader.pages)
    results = []
    for start, end in ranges:
        start = max(1, start)
        end = min(total, end)
        writer = PdfWriter()
        for i in range(start - 1, end):
            writer.add_page(reader.pages[i])
        buf = io.BytesIO()
        writer.write(buf)
        results.append(buf.getvalue())
    return results


def compress_pdf(data: bytes) -> bytes:
    """
    Lossless-ish compression via pikepdf: object-stream compression, removes
    unused objects, and re-encodes. For aggressive image downsampling you'd
    add Ghostscript (see README) — this is the safe, dependency-light path.
    """
    with pikepdf.open(io.BytesIO(data)) as pdf:
        out = io.BytesIO()
        pdf.save(
            out,
            compress_streams=True,
            object_stream_mode=pikepdf.ObjectStreamMode.generate,
            recompress_flate=True,
        )
        return out.getvalue()


def rotate_pdf(data: bytes, degrees: int = 90) -> bytes:
    """Rotate every page by `degrees` (must be a multiple of 90)."""
    reader = PdfReader(io.BytesIO(data))
    writer = PdfWriter()
    for page in reader.pages:
        page.rotate(degrees)
        writer.add_page(page)
    out = io.BytesIO()
    writer.write(out)
    return out.getvalue()


def pdf_page_count(data: bytes) -> int:
    return len(PdfReader(io.BytesIO(data)).pages)


def protect_pdf(data: bytes, password: str) -> bytes:
    """Encrypt a PDF with a password (AES-256)."""
    with pikepdf.open(io.BytesIO(data)) as pdf:
        out = io.BytesIO()
        pdf.save(
            out,
            encryption=pikepdf.Encryption(owner=password, user=password, R=6),
        )
        return out.getvalue()


def unlock_pdf(data: bytes, password: str) -> bytes:
    """
    Remove a known password from a PDF. Raises pikepdf.PasswordError if the
    password is wrong (the view turns that into a clean 400).
    """
    with pikepdf.open(io.BytesIO(data), password=password) as pdf:
        out = io.BytesIO()
        pdf.save(out)  # saving without an encryption arg drops the encryption
        return out.getvalue()


def repair_pdf(data: bytes) -> bytes:
    """
    Best-effort repair: pikepdf re-parses and rewrites the file, rebuilding the
    cross-reference table and cleaning up structural damage. Truly corrupt files
    may still fail (the view returns a clean error then).
    """
    with pikepdf.open(io.BytesIO(data)) as pdf:
        out = io.BytesIO()
        pdf.save(out, fix_metadata_version=True)
        return out.getvalue()


def pdf_to_images(data: bytes, dpi: int = 150, fmt: str = "jpeg") -> List[bytes]:
    """
    Render each PDF page to a raster image using PyMuPDF (fitz).
    We use PyMuPDF instead of pdf2image/poppler on purpose: it ships as a
    self-contained wheel, so there are NO system packages to install (works
    the same on Windows dev and in the slim Docker image).
    Returns one image (bytes) per page.
    """
    import fitz  # PyMuPDF — imported lazily so the other tools work without it

    images: List[bytes] = []
    with fitz.open(stream=data, filetype="pdf") as doc:
        for page in doc:
            pix = page.get_pixmap(dpi=dpi)
            images.append(pix.tobytes(fmt))
    return images


def pdf_to_docx(data: bytes) -> bytes:
    """
    Convert a PDF into an editable Word (.docx) using pdf2docx (pure pip, no
    system binaries). It reconstructs text, tables and basic layout.
    """
    import os
    import tempfile

    from pdf2docx import Converter  # lazy import — keeps other tools dep-free

    with tempfile.TemporaryDirectory() as tmp:
        pdf_path = os.path.join(tmp, "in.pdf")
        docx_path = os.path.join(tmp, "out.docx")
        with open(pdf_path, "wb") as fh:
            fh.write(data)
        cv = Converter(pdf_path)
        try:
            cv.convert(docx_path)  # all pages
        finally:
            cv.close()
        with open(docx_path, "rb") as fh:
            return fh.read()


def docx_to_pdf(data: bytes, filename: str = "document.docx") -> bytes:
    """
    Convert a Word document to PDF via LibreOffice headless (`soffice`).
    This is the one tool that needs a system binary. If LibreOffice isn't
    installed we raise a clear error (the view turns it into a 503) instead of
    crashing — so this tool degrades gracefully and never affects the others.
    """
    import glob
    import os
    import shutil
    import subprocess
    import tempfile

    # Find LibreOffice. On PATH (Docker/Linux) or common Windows install dirs,
    # since the Windows installer doesn't add `soffice` to PATH.
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        for candidate in (
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        ):
            if os.path.exists(candidate):
                soffice = candidate
                break
    if not soffice:
        raise RuntimeError(
            "This conversion needs LibreOffice on the server. Install the "
            "'libreoffice' package (already added to the Dockerfile) and "
            "redeploy, or install LibreOffice locally to test on Windows."
        )

    with tempfile.TemporaryDirectory() as tmp:
        ext = os.path.splitext(filename)[1].lower() or ".docx"
        in_path = os.path.join(tmp, "input" + ext)
        with open(in_path, "wb") as fh:
            fh.write(data)
        # A private user-profile dir avoids clashes when requests run in parallel.
        profile = "-env:UserInstallation=file://" + os.path.join(tmp, "lo_profile")
        subprocess.run(
            [soffice, profile, "--headless", "--convert-to", "pdf",
             "--outdir", tmp, in_path],
            check=True, timeout=120,
            stdout=subprocess.PIPE, stderr=subprocess.PIPE,
        )
        pdfs = glob.glob(os.path.join(tmp, "*.pdf"))
        if not pdfs:
            raise RuntimeError("LibreOffice produced no PDF output.")
        with open(pdfs[0], "rb") as fh:
            return fh.read()


# LibreOffice converts any office/HTML format to PDF purely from the file
# extension, so PowerPoint / Excel / HTML → PDF all reuse the same routine.
office_to_pdf = docx_to_pdf


def pdf_to_markdown(data: bytes) -> bytes:
    """Extract a PDF's text into Markdown (one section per page)."""
    import fitz

    parts = []
    with fitz.open(stream=data, filetype="pdf") as doc:
        for i, page in enumerate(doc, start=1):
            text = page.get_text("text").strip()
            parts.append(f"## Page {i}\n\n{text}\n")
    md = "\n---\n\n".join(parts).strip()
    if not md:
        md = ("# No selectable text found\n\n"
              "This PDF looks scanned — run it through **OCR PDF** first, then "
              "convert to Markdown.")
    return md.encode("utf-8")


def pdf_to_pptx(data: bytes, dpi: int = 150) -> bytes:
    """Render each PDF page to an image and drop it onto its own slide."""
    import io

    import fitz
    from pptx import Presentation
    from pptx.util import Emu

    pages = []
    with fitz.open(stream=data, filetype="pdf") as doc:
        for page in doc:
            pix = page.get_pixmap(dpi=dpi)
            pages.append((pix.tobytes("png"), page.rect.width, page.rect.height))

    prs = Presentation()
    if pages:
        # PowerPoint uses one size for the whole deck; take the first page (pt→EMU).
        _, w_pt, h_pt = pages[0]
        prs.slide_width = Emu(int(w_pt * 12700))
        prs.slide_height = Emu(int(h_pt * 12700))
    blank = prs.slide_layouts[6]
    for png, _w, _h in pages:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            io.BytesIO(png), 0, 0, width=prs.slide_width, height=prs.slide_height
        )
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def pdf_to_xlsx(data: bytes) -> bytes:
    """Pull tables (or failing that, text lines) from a PDF into a spreadsheet."""
    import io

    import pdfplumber
    from openpyxl import Workbook

    wb = Workbook()
    wb.remove(wb.active)
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for pi, page in enumerate(pdf.pages, start=1):
            ws = wb.create_sheet(title=f"Page {pi}"[:31])
            r = 1
            tables = page.extract_tables()
            if tables:
                for table in tables:
                    for row in table:
                        for ci, val in enumerate(row, start=1):
                            ws.cell(row=r, column=ci, value=(val or ""))
                        r += 1
                    r += 1  # blank row between tables
            else:
                for line in (page.extract_text() or "").splitlines():
                    ws.cell(row=r, column=1, value=line)
                    r += 1
    if not wb.sheetnames:
        wb.create_sheet(title="Empty")
    out = io.BytesIO()
    wb.save(out)
    return out.getvalue()


def ocr_pdf(data: bytes, lang: str = "eng") -> bytes:
    """
    Add a searchable text layer to a scanned PDF via OCRmyPDF (Tesseract).
    Needs the tesseract-ocr + ghostscript binaries — raises RuntimeError (→ 503)
    when they're missing so the tool degrades gracefully.
    """
    import os
    import tempfile

    try:
        import ocrmypdf
    except Exception as e:  # package not installed
        raise RuntimeError(
            "OCR needs 'ocrmypdf' + Tesseract on the server (added to the "
            "Dockerfile). Redeploy the backend to enable it."
        ) from e

    with tempfile.TemporaryDirectory() as tmp:
        inp = os.path.join(tmp, "in.pdf")
        outp = os.path.join(tmp, "out.pdf")
        with open(inp, "wb") as fh:
            fh.write(data)
        try:
            ocrmypdf.ocr(
                inp, outp, language=lang, skip_text=True, progress_bar=False,
            )
        except Exception as e:
            raise RuntimeError(
                "OCR could not run — Tesseract/Ghostscript may be missing on the "
                "server. Redeploy with the updated Dockerfile."
            ) from e
        with open(outp, "rb") as fh:
            return fh.read()


def _extract_text(data: bytes, max_chars: int = 60000) -> str:
    """Pull a PDF's text (capped) for the AI tools to work on."""
    import fitz

    parts, total = [], 0
    with fitz.open(stream=data, filetype="pdf") as doc:
        for page in doc:
            t = page.get_text("text")
            parts.append(t)
            total += len(t)
            if total >= max_chars:
                break
    return "\n".join(parts).strip()[:max_chars]


def _claude(prompt: str, max_tokens: int) -> str:
    """
    Call Claude. Requires ANTHROPIC_API_KEY on the server — raises RuntimeError
    (→ 503) when it's absent so the AI tools degrade gracefully. Model defaults
    to claude-opus-5; override with ANTHROPIC_MODEL (e.g. claude-haiku-4-5 to
    cut cost).
    """
    import os

    if not os.environ.get("ANTHROPIC_API_KEY"):
        raise RuntimeError(
            "AI features need an Anthropic API key. Set ANTHROPIC_API_KEY on the "
            "server (and optionally ANTHROPIC_MODEL) and redeploy."
        )
    import anthropic

    model = os.environ.get("ANTHROPIC_MODEL", "claude-opus-5")
    client = anthropic.Anthropic()
    with client.messages.stream(
        model=model,
        max_tokens=max_tokens,
        messages=[{"role": "user", "content": prompt}],
    ) as stream:
        msg = stream.get_final_message()
    return "".join(b.text for b in msg.content if b.type == "text").strip()


def summarize_pdf(data: bytes) -> bytes:
    text = _extract_text(data)
    if not text:
        raise ValueError("No selectable text found — run OCR first, then summarize.")
    prompt = (
        "Summarize the following document. Start with a short list of the key "
        "points as bullets, then a one-paragraph overview. Be clear and concise.\n\n"
        "---\n\n" + text
    )
    return _claude(prompt, max_tokens=4096).encode("utf-8")


def translate_pdf(data: bytes, target_language: str) -> bytes:
    text = _extract_text(data)
    if not text:
        raise ValueError("No selectable text found — run OCR first, then translate.")
    prompt = (
        f"Translate the following document into {target_language}. Preserve the "
        "structure and meaning. Output only the translation, with no commentary.\n\n"
        "---\n\n" + text
    )
    return _claude(prompt, max_tokens=16000).encode("utf-8")


def compare_pdfs(a: bytes, b: bytes, dpi: int = 110, quality: int = 78) -> bytes:
    """
    Compare two PDFs page by page and produce a side-by-side report: the
    original on the left, the changed version on the right, with the regions
    that differ highlighted on both.

    WHY SIDE-BY-SIDE RATHER THAN AN OVERLAY
    The previous version composited one page on top of the other and painted
    every differing pixel OPAQUE red. That only reads correctly when the two
    files are already nearly identical. Give it two genuinely different
    documents — the common case, and what a real catalogue produced — and
    almost every pixel differs, so the page came back a solid red block with
    both documents' content smeared through it. Measured on dense pages: 42-48%
    of the sheet painted flat red, and the underlying content destroyed.

    The fix is twofold. Showing both pages next to each other means you can
    actually read what changed instead of inferring it, and the highlight is
    TRANSLUCENT, so content stays legible underneath rather than being replaced
    by a colour block.

    Also kept from the previous fix: JPEG rather than lossless PNG (a 2-page
    diff was 8MB as PNG), and page geometry converted from pixels to points, so
    output prints at 1:1 instead of dpi/72 oversized.
    """
    import io

    import fitz
    from PIL import Image, ImageChops, ImageFilter

    GAP = 16.0        # points between the two pages
    MARGIN = 14.0
    HEADER = 30.0     # room for the labels above each page

    def render(doc, i):
        if i >= len(doc):
            return None
        pix = doc[i].get_pixmap(dpi=dpi)
        return Image.frombytes("RGB", (pix.width, pix.height), pix.samples)

    def encode(img: Image.Image) -> bytes:
        buf = io.BytesIO()
        img.save(buf, format="JPEG", quality=quality, optimize=True)
        return buf.getvalue()

    def tint(img: Image.Image, mask: Image.Image, colour: tuple[int, int, int]) -> Image.Image:
        """Wash the masked regions with colour, keeping the content readable."""
        out = img.copy()
        layer = Image.new("RGB", img.size, colour)
        # ~40% opacity. Enough to spot at a glance, light enough to read through.
        out.paste(layer, (0, 0), mask.point(lambda p: 105 if p else 0))
        return out

    da = fitz.open(stream=a, filetype="pdf")
    db = fitz.open(stream=b, filetype="pdf")
    try:
        out = fitz.open()
        n = max(len(da), len(db))

        for i in range(n):
            ia = render(da, i)
            ib = render(db, i)

            changed_pct = None
            if ia is not None and ib is not None:
                w = max(ia.width, ib.width)
                h = max(ia.height, ib.height)
                # Pad to a common canvas so pages of differing size still align
                # at the top-left rather than failing to subtract.
                pa = Image.new("RGB", (w, h), "white")
                pb = Image.new("RGB", (w, h), "white")
                pa.paste(ia, (0, 0))
                pb.paste(ib, (0, 0))

                diff = ImageChops.difference(pa, pb).convert("L")
                mask = diff.point(lambda p: 255 if p > 32 else 0)
                # Consolidate speckle from anti-aliasing into solid regions, so
                # the highlight marks areas rather than peppering the page.
                mask = mask.filter(ImageFilter.MaxFilter(5))

                changed_pct = (sum(mask.point(lambda p: 1 if p else 0).getdata()) / (w * h)) * 100.0

                left = tint(pa, mask, (220, 60, 60))    # removed / was here
                right = tint(pb, mask, (70, 170, 90))   # added / is here now
            else:
                # Page exists in only one document.
                left = ia
                right = ib

            # ---- compose the output page -------------------------------------
            def pt(img):
                return (img.width * 72.0 / dpi, img.height * 72.0 / dpi) if img else (0.0, 0.0)

            lw, lh = pt(left)
            rw, rh = pt(right)
            page_w = MARGIN * 2 + lw + GAP + rw
            page_h = MARGIN * 2 + HEADER + max(lh, rh)
            page = out.new_page(width=page_w, height=page_h)

            y0 = MARGIN + HEADER
            if left is not None:
                page.insert_image(fitz.Rect(MARGIN, y0, MARGIN + lw, y0 + lh), stream=encode(left))
                page.insert_text(
                    (MARGIN, MARGIN + 14),
                    f"Original — page {i + 1}" if ia is not None else "",
                    fontname="hebo", fontsize=10, color=(0.35, 0.35, 0.4),
                )
            if right is not None:
                x = MARGIN + lw + GAP
                page.insert_image(fitz.Rect(x, y0, x + rw, y0 + rh), stream=encode(right))
                page.insert_text(
                    (x, MARGIN + 14),
                    f"Changed — page {i + 1}" if ib is not None else "",
                    fontname="hebo", fontsize=10, color=(0.35, 0.35, 0.4),
                )

            # A one-line verdict beats making the reader estimate from colour.
            if changed_pct is not None:
                note = (
                    "identical" if changed_pct < 0.05
                    else f"{changed_pct:.1f}% of the page differs"
                )
                colour = (0.2, 0.55, 0.3) if changed_pct < 0.05 else (0.75, 0.2, 0.2)
                page.insert_text(
                    (MARGIN, MARGIN + 27), note,
                    fontname="helv", fontsize=8, color=colour,
                )
            else:
                only = "original" if ib is None else "changed version"
                page.insert_text(
                    (MARGIN, MARGIN + 27),
                    f"page {i + 1} exists only in the {only}",
                    fontname="helv", fontsize=8, color=(0.75, 0.2, 0.2),
                )

        return out.tobytes(deflate=True, garbage=3)
    finally:
        da.close()
        db.close()


def pdf_to_pdfa(data: bytes) -> bytes:
    """
    Convert to PDF/A-2b for long-term archiving via Ghostscript. Raises
    RuntimeError (→ 503) if Ghostscript isn't installed.
    """
    import os
    import shutil
    import subprocess
    import tempfile

    gs = shutil.which("gs") or shutil.which("gswin64c") or shutil.which("gswin32c")
    if not gs:
        raise RuntimeError(
            "PDF/A needs Ghostscript on the server (added to the Dockerfile). "
            "Redeploy the backend to enable it."
        )
    with tempfile.TemporaryDirectory() as tmp:
        inp = os.path.join(tmp, "in.pdf")
        outp = os.path.join(tmp, "out.pdf")
        with open(inp, "wb") as fh:
            fh.write(data)
        subprocess.run(
            [gs, "-dPDFA=2", "-dBATCH", "-dNOPAUSE", "-dQUIET",
             "-sColorConversionStrategy=RGB", "-sDEVICE=pdfwrite",
             "-dPDFACompatibilityPolicy=1", "-sOutputFile=" + outp, inp],
            check=True, timeout=180,
            stdout=subprocess.PIPE, stderr=subprocess.PIPE,
        )
        with open(outp, "rb") as fh:
            return fh.read()
